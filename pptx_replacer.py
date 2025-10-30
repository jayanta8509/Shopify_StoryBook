"""
PowerPoint Text Replacer
Replace placeholder text like {{CHILD_NAME}} in .pptx files
"""

from pptx import Presentation
from pathlib import Path
import re
from typing import Dict, List
import copy


class PowerPointReplacer:
    def __init__(self, template_path: str):
        """
        Initialize the PowerPoint Replacer
        
        Args:
            template_path: Path to the template .pptx file
        """
        self.template_path = Path(template_path)
        if not self.template_path.exists():
            raise FileNotFoundError(f"Template file not found: {template_path}")
    
    def find_placeholders(self) -> List[str]:
        """
        Find all unique placeholders in the presentation
        
        Returns:
            List of unique placeholder strings (e.g., ['{{CHILD_NAME}}', '{{CHILD_NAME_UPPER}}'])
        """
        prs = Presentation(self.template_path)
        placeholders = set()
        
        # Pattern to match placeholders like {{SOMETHING}}
        pattern = re.compile(r'\{\{[A-Z_]+\}\}')
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        matches = pattern.findall(run.text)
                        placeholders.update(matches)
        
        return sorted(list(placeholders))
    
    def replace_text(self, replacements: Dict[str, str], output_path: str) -> str:
        """
        Replace text in the PowerPoint file
        
        Args:
            replacements: Dictionary of {placeholder: replacement_text}
                         e.g., {'{{CHILD_NAME}}': 'Emma', '{{CHILD_NAME_UPPER}}': 'EMMA'}
            output_path: Path where the modified file should be saved
        
        Returns:
            Path to the output file
        """
        # Load the presentation
        prs = Presentation(self.template_path)
        
        # Track statistics
        total_replacements = 0
        slides_modified = 0
        
        # Process each slide
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_modified = False
            
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                
                # Process each paragraph in the text frame
                for paragraph in shape.text_frame.paragraphs:
                    # Process each run in the paragraph
                    for run in paragraph.runs:
                        original_text = run.text
                        modified_text = original_text
                        
                        # Apply all replacements
                        for placeholder, replacement in replacements.items():
                            if placeholder in modified_text:
                                modified_text = modified_text.replace(placeholder, replacement)
                                total_replacements += 1
                                slide_modified = True
                        
                        # Update the run text if it was modified
                        if modified_text != original_text:
                            run.text = modified_text
            
            if slide_modified:
                slides_modified += 1
        
        # Save the modified presentation
        output_path = Path(output_path)
        prs.save(str(output_path))
        
        print(f"✅ Successfully created personalized presentation!")
        print(f"   - Total replacements: {total_replacements}")
        print(f"   - Slides modified: {slides_modified}/{len(prs.slides)}")
        print(f"   - Saved to: {output_path}")
        
        return str(output_path)
    
    def create_multiple(
        self, 
        names_list: List[Dict[str, str]], 
        output_dir: str,
        filename_pattern: str = "{name}_personalized.pptx"
    ) -> List[str]:
        """
        Create multiple personalized presentations
        
        Args:
            names_list: List of dictionaries with placeholder mappings
                       e.g., [
                           {'{{CHILD_NAME}}': 'Emma', '{{CHILD_NAME_UPPER}}': 'EMMA'},
                           {'{{CHILD_NAME}}': 'Liam', '{{CHILD_NAME_UPPER}}': 'LIAM'}
                       ]
            output_dir: Directory to save all output files
            filename_pattern: Pattern for output filenames (use {name} for child name)
        
        Returns:
            List of paths to created files
        """
        output_dir = Path(output_dir)
        output_dir.mkdir(parents=True, exist_ok=True)
        
        created_files = []
        
        for i, replacements in enumerate(names_list, 1):
            # Extract the child name for the filename
            child_name = replacements.get('{{CHILD_NAME}}', f'child_{i}')
            
            # Create output filename
            output_filename = filename_pattern.format(name=child_name)
            output_path = output_dir / output_filename
            
            print(f"\n📄 Creating presentation {i}/{len(names_list)} for {child_name}...")
            
            # Replace text and save
            created_file = self.replace_text(replacements, str(output_path))
            created_files.append(created_file)
        
        print(f"\n🎉 Successfully created {len(created_files)} personalized presentations!")
        
        return created_files


def main():
    """Example usage"""
    
    # Path to your template file
    template_path = 'Storybook_Template_DK BOY_v1.pptx'
    
    # Create replacer instance
    replacer = PowerPointReplacer(template_path)
    
    # Option 1: Find what placeholders exist in the template
    print("🔍 Finding placeholders in template...")
    placeholders = replacer.find_placeholders()
    print(f"Found placeholders: {placeholders}")
    
    # Option 2: Create a single personalized presentation
    print("\n" + "="*60)
    print("Creating single personalized presentation...")
    print("="*60)
    
    replacements = {
        '{{CHILD_NAME}}': 'Emma',
        '{{CHILD_NAME_UPPER}}': 'EMMA'
    }
    
    output_path = 'media/Emma_Storybook.pptx'
    replacer.replace_text(replacements, output_path)
    
    # Option 3: Create multiple personalized presentations
    print("\n" + "="*60)
    print("Creating multiple personalized presentations...")
    print("="*60)
    
    children = [
        {'{{CHILD_NAME}}': 'Liam', '{{CHILD_NAME_UPPER}}': 'LIAM'},
        {'{{CHILD_NAME}}': 'Olivia', '{{CHILD_NAME_UPPER}}': 'OLIVIA'},
        {'{{CHILD_NAME}}': 'Noah', '{{CHILD_NAME_UPPER}}': 'NOAH'},
    ]
    
    output_dir = 'media'
    created_files = replacer.create_multiple(
        children, 
        output_dir,
        filename_pattern="{name}_Storybook.pptx"
    )
    
    print("\n✨ All done! Created files:")
    for file_path in created_files:
        print(f"   - {file_path}")


if __name__ == '__main__':
    main()
